"""
Bettroi — AI-Powered Retail Intelligence Platform
PowerPoint Deck Generator for Ajmal Perfumes Pitch
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BRAND_PRIMARY = RGBColor(0x6C, 0x63, 0xFF)
BRAND_SECONDARY = RGBColor(0x00, 0xD4, 0xAA)
BRAND_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)
BRAND_GOLD = RGBColor(0xF5, 0xA6, 0x23)
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x16, 0x16, 0x22)
TEXT_PRIMARY = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_SECONDARY = RGBColor(0x88, 0x88, 0xAA)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=TEXT_SECONDARY,
                  bold=False, space_before=0, space_after=0, alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = alignment
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.6, 4, 0.5, 'BETTROI', font_size=28, color=WHITE, bold=True)
add_textbox(slide, 3.2, 0.7, 2, 0.4, 'est. 2024', font_size=12, color=TEXT_MUTED)

add_textbox(slide, 0.8, 1.8, 10, 1.2, 'AI-Powered', font_size=52, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 2.8, 10, 1.2, 'Retail Intelligence', font_size=52, color=WHITE, bold=True)
add_textbox(slide, 0.8, 3.8, 10, 1.2, 'Platform', font_size=52, color=WHITE, bold=True)

add_textbox(slide, 0.8, 5.2, 8, 0.5, "Built for India's Premium Retail Chains", font_size=22, color=TEXT_SECONDARY)

add_textbox(slide, 0.8, 6.4, 5, 0.4, 'FROM BILLING TO BRILLIANCE', font_size=11, color=TEXT_MUTED, bold=True)

add_notes(slide, "Keep this slide up during introductions. Let the tagline 'From Billing to Brilliance' sit — it plants the POS-to-AI journey in their mind from the start. Don't explain it yet.")


# ============================================================
# SLIDE 2: INDUSTRY LANDSCAPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'INDUSTRY LANDSCAPE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, "India's Organized Retail Is at an Inflection Point", font_size=36, color=WHITE, bold=True)

# Stat boxes
stats = [
    ("$28B", "India Beauty & Personal Care\nMarket by 2028"),
    ("18-22%", "Premium Fragrance\nCAGR"),
    ("78%", "Retail Chains Still on\nLegacy Billing Systems"),
    ("2.3x", "Faster Growth for Chains\nwith Real-Time Analytics"),
]

for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 3.05
    add_rounded_rect(slide, x, 2.2, 2.8, 2.0, BG_CARD)
    add_textbox(slide, x, 2.4, 2.8, 0.8, num, font_size=40, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 3.4, 2.5, 0.8, label, font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 4.8, 11, 1.0,
    "The industry is growing fast — but most retail infrastructure was built for a different era. "
    "The winners will be the chains that close the gap between brand ambition and operational backbone.",
    font_size=16, color=TEXT_SECONDARY)

add_notes(slide, "Let the '78% legacy billing' stat land. Don't point at anyone — let them self-identify. 'Retailers with real-time analytics grow 2.3x faster' creates urgency. The audience should be thinking: 'Are we in the 78% or the 22%?'")


# ============================================================
# SLIDE 3: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'THE PATTERN WE SEE', font_size=11, color=BRAND_ACCENT, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'What We Hear From Every Growing Retail Chain', font_size=36, color=WHITE, bold=True)

# Brand side
add_rounded_rect(slide, 0.8, 2.2, 5.5, 2.8, BG_CARD, RGBColor(0x22, 0x44, 0x22))
add_textbox(slide, 1.0, 2.3, 5, 0.4, 'The Brand Side  ✓', font_size=16, color=BRAND_SECONDARY, bold=True)
brand_items = [
    "→  Beautiful stores, strong heritage",
    "→  Loyal customer base, growing footprint",
    "→  Premium positioning, expanding portfolio",
    "→  Multi-channel presence (retail, e-commerce, shop-in-shop)"
]
for j, item in enumerate(brand_items):
    add_textbox(slide, 1.0, 2.8 + j * 0.42, 5, 0.4, item, font_size=13, color=TEXT_SECONDARY)

# Operations side
add_rounded_rect(slide, 6.8, 2.2, 5.7, 2.8, BG_CARD, RGBColor(0x44, 0x22, 0x22))
add_textbox(slide, 7.0, 2.3, 5, 0.4, 'The Operations Side  ⚠', font_size=16, color=BRAND_ACCENT, bold=True)
ops_items = [
    "→  Manual billing slowing down peak hours",
    "→  No real-time visibility across stores",
    "→  HQ decisions based on yesterday's data",
    "→  Staff performance invisible until quarterly reviews",
    "→  Franchise and owned stores on different systems"
]
for j, item in enumerate(ops_items):
    add_textbox(slide, 7.0, 2.8 + j * 0.42, 5.3, 0.4, item, font_size=13, color=BRAND_ACCENT)

# Quote
add_rounded_rect(slide, 0.8, 5.3, 11.7, 1.0, RGBColor(0x12, 0x12, 0x20), BRAND_PRIMARY)
add_textbox(slide, 1.2, 5.4, 10.5, 0.5,
    '"We have 100+ stores but our HQ team spends Monday mornings calling store managers for weekend numbers."',
    font_size=14, color=TEXT_SECONDARY)
add_textbox(slide, 1.2, 5.9, 10, 0.3, '— Operations Head, Premium Retail Chain', font_size=11, color=TEXT_MUTED)

add_textbox(slide, 0.8, 6.7, 10, 0.4, 'Tools that worked at 30 stores break at 100. And they completely fail at 150.',
    font_size=12, color=TEXT_MUTED)

add_notes(slide, "This is the 'we see you' slide. Frame as industry-wide, NOT Ajmal-specific. 'We hear some version of this in every initial conversation.' The anonymous quote should feel eerily familiar. The franchise/owned distinction is deliberate — they run both.")


# ============================================================
# SLIDE 4: PLATFORM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'PLATFORM OVERVIEW', font_size=11, color=BRAND_SECONDARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'One Platform. Every Store. Every Insight.', font_size=36, color=WHITE, bold=True)

# Experience Layer
add_rounded_rect(slide, 0.8, 2.2, 11.7, 1.3, RGBColor(0x14, 0x12, 0x28), BRAND_PRIMARY)
add_textbox(slide, 1.0, 2.3, 4, 0.3, 'EXPERIENCE LAYER', font_size=10, color=BRAND_PRIMARY, bold=True)
exp_items = ['AI Fragrance Advisor', 'WhatsApp CRM', 'Dynamic Pricing', 'L&D Academy']
for i, item in enumerate(exp_items):
    x = 0.9 + i * 2.9
    add_rounded_rect(slide, x, 2.7, 2.6, 0.5, RGBColor(0x1A, 0x18, 0x30))
    add_textbox(slide, x, 2.75, 2.6, 0.4, item, font_size=12, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

# Intelligence Layer
add_rounded_rect(slide, 0.8, 3.7, 11.7, 1.3, RGBColor(0x0E, 0x1A, 0x18), BRAND_SECONDARY)
add_textbox(slide, 1.0, 3.8, 4, 0.3, 'INTELLIGENCE LAYER', font_size=10, color=BRAND_SECONDARY, bold=True)
int_items = ['Demand Forecasting', 'Customer Analytics', 'Inventory Intelligence', 'Staff Performance AI']
for i, item in enumerate(int_items):
    x = 1.0 + i * 2.9
    add_rounded_rect(slide, x, 4.2, 2.6, 0.5, RGBColor(0x12, 0x20, 0x1C))
    add_textbox(slide, x, 4.25, 2.6, 0.4, item, font_size=12, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

# Foundation Layer (highlighted)
add_rounded_rect(slide, 0.8, 5.2, 11.7, 1.3, RGBColor(0x1E, 0x16, 0x0A), BRAND_GOLD)
add_textbox(slide, 1.0, 5.3, 6, 0.3, '⬤  FOUNDATION LAYER — START HERE', font_size=10, color=BRAND_GOLD, bold=True)
fnd_items = ['Smart POS & Store Operations', 'Real-Time Command Center & Dashboards']
for i, item in enumerate(fnd_items):
    x = 1.2 + i * 5.5
    add_rounded_rect(slide, x, 5.7, 5.0, 0.5, RGBColor(0x28, 0x20, 0x10), BRAND_GOLD)
    add_textbox(slide, x, 5.75, 5.0, 0.4, item, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 6.7, 11.7, 0.4,
    'The foundation powers the intelligence. The intelligence powers the experience.',
    font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "THIS is the key architectural slide. POS and dashboards are the mandatory foundation — not because we're pushing them, but because AI layers cannot work without clean operational data. 'You cannot run AI on bad data' is unassailable logic. The golden border on the foundation draws the eye naturally.")


# ============================================================
# SLIDE 5: SMART POS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'FOUNDATION MODULE', font_size=11, color=BRAND_GOLD, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Smart POS — The Nerve Center of Every Store', font_size=36, color=WHITE, bold=True)

# Features
features_pos = [
    "→  One-tap billing with auto product recognition",
    "→  Auto-applied promotions — no manual discount entry",
    "→  Offline-first — works without internet, syncs when back",
    "→  Franchise + company-owned on one unified system",
    "→  Multi-format: full store, kiosk, shop-in-shop counter",
    "→  Integrated payments: UPI, card, wallet — auto-reconciled",
    "→  GST-ready invoicing with instant digital receipts",
]
for j, feat in enumerate(features_pos):
    add_textbox(slide, 0.8, 2.1 + j * 0.42, 6, 0.4, feat, font_size=13, color=TEXT_SECONDARY)

# Metrics
add_rounded_rect(slide, 7.5, 2.1, 5.0, 3.5, BG_CARD)
add_textbox(slide, 7.7, 2.2, 4.5, 0.4, 'What Changes', font_size=18, color=BRAND_SECONDARY, bold=True)

metrics = [
    ("4-5 min", "→", "45 sec", "Billing time per transaction"),
    ("Manual", "→", "Auto", "Promo & discount application"),
    ("90 min", "→", "5 min", "End-of-day reconciliation"),
    ("Next day", "→", "Real-time", "Data available to HQ"),
]
for j, (before, arrow, after, label) in enumerate(metrics):
    y = 2.9 + j * 0.7
    add_textbox(slide, 7.7, y, 1.3, 0.3, before, font_size=13, color=TEXT_MUTED)
    add_textbox(slide, 9.1, y, 0.4, 0.3, arrow, font_size=13, color=BRAND_SECONDARY, bold=True)
    add_textbox(slide, 9.5, y, 1.3, 0.3, after, font_size=13, color=BRAND_SECONDARY, bold=True)
    add_textbox(slide, 7.7, y + 0.25, 4.5, 0.3, label, font_size=11, color=TEXT_MUTED)

add_textbox(slide, 0.8, 6.5, 11, 0.4, 'Every transaction becomes a data point. Every data point powers AI.',
    font_size=12, color=TEXT_MUTED)

add_notes(slide, "'Most POS systems were designed as cash registers that happen to be digital. Ours was designed as a data capture engine that happens to do billing.' Offline-first matters for India — mall WiFi drops during peak hours. Franchise standardization is built in — HQ gets one clean data feed, not 50 Excel formats.")


# ============================================================
# SLIDE 6: COMMAND CENTER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'FOUNDATION MODULE', font_size=11, color=BRAND_GOLD, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Command Center — See Everything. Right Now.', font_size=36, color=WHITE, bold=True)

# Live Metrics
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.0, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'LIVE METRICS', font_size=11, color=BRAND_SECONDARY, bold=True)
metrics_list = [
    "→  Store-by-store sales — today, this week, vs target",
    "→  Staff performance heatmap by store & region",
    "→  Top-selling SKUs by city, format, season",
    "→  Franchise vs. company-owned comparison",
    "→  Alert flags: below target, stock-out, high returns",
]
for j, m in enumerate(metrics_list):
    add_textbox(slide, 1.0, 2.7 + j * 0.42, 5, 0.4, m, font_size=13, color=TEXT_SECONDARY)

# Key Capabilities
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.0, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'KEY CAPABILITIES', font_size=11, color=BRAND_GOLD, bold=True)
caps = [
    "→  Zero manual data entry — feeds directly from POS",
    "→  Role-based views: store mgr, area head, HQ, CEO",
    "→  Mobile-first: CEO checks Sunday sales from Dubai\n     on Monday morning",
    "→  Auto-reports: daily, weekly, monthly — no spreadsheets",
    "→  Franchise owner portal with controlled visibility",
]
for j, c in enumerate(caps):
    add_textbox(slide, 7.0, 2.7 + j * 0.45, 5.3, 0.45, c, font_size=13, color=TEXT_SECONDARY)

# Quote
add_rounded_rect(slide, 0.8, 5.4, 11.7, 0.9, RGBColor(0x12, 0x12, 0x20), BRAND_PRIMARY)
add_textbox(slide, 1.2, 5.5, 10.5, 0.6,
    'The #1 complaint from HQ teams: "I cannot get accurate store data without calling someone." This eliminates that — entirely.',
    font_size=14, color=TEXT_SECONDARY)

add_notes(slide, "Say 'CEO checks Sunday sales from Dubai on Monday morning' casually — as if this is a common scenario for distributed retail chains. The 'zero manual data entry' line directly addresses their pain, framed as universal. The franchise owner portal shows you understand hybrid models.")


# ============================================================
# SLIDE 7: AI FRAGRANCE ADVISOR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'AI Fragrance Advisor — Turn Every Associate Into an Expert', font_size=34, color=WHITE, bold=True)

# 3-step flow
steps = [
    ("1", "Customer shares\npreferences", BRAND_PRIMARY),
    ("2", "AI matches fragrance\nknowledge graph", BRAND_SECONDARY),
    ("3", "3-4 recommendations\nwith talking points", BRAND_GOLD),
]
for i, (num, label, color) in enumerate(steps):
    x = 1.5 + i * 4.0
    add_rounded_rect(slide, x, 2.1, 0.6, 0.6, BG_CARD, color)
    add_textbox(slide, x, 2.15, 0.6, 0.5, num, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.8, 2.1, 2.5, 0.6, label, font_size=13, color=TEXT_SECONDARY)

# Features
features_frag = [
    "→  Matches by occasion, intensity, budget, note preferences",
    "→  Learns from purchase history: 'Customers who bought X love Y'",
    "→  Cross-sell: attar buyer → suggest matching bakhoor",
    "→  Handles full range: ₹375 body mist to ₹7,500+ luxury oud",
    "→  Works on store tablet or associate's phone",
]
for j, f in enumerate(features_frag):
    add_textbox(slide, 0.8, 3.2 + j * 0.42, 7, 0.4, f, font_size=13, color=TEXT_SECONDARY)

# Stat
add_rounded_rect(slide, 8.5, 3.2, 4.0, 2.0, BG_CARD)
add_textbox(slide, 8.5, 3.5, 4.0, 0.8, '22-35%', font_size=44, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 8.5, 4.3, 4.0, 0.6, 'Increase in average\nbasket size', font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 5.8, 10, 0.5,
    "Your best salesperson has 10 years of fragrance knowledge. This gives a new hire 80% of that on day one.",
    font_size=15, color=TEXT_PRIMARY)

add_notes(slide, "This is the WOW slide. 'Fragrance is one of the hardest categories — nose fatigue after 3-4 tries. AI narrows the field.' The cross-sell from EDP to bakhoor/attar is where real margin uplift happens. No other retail tech platform has a fragrance-specific recommendation engine.")


# ============================================================
# SLIDE 8: DEMAND & INVENTORY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Know What Will Sell, Where, Before It Happens', font_size=36, color=WHITE, bold=True)

# Demand Forecasting
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.2, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'DEMAND FORECASTING', font_size=11, color=BRAND_SECONDARY, bold=True)
demand_items = [
    "→  AI predicts SKU-level demand by store, by week",
    "→  Factors: Eid, Diwali, wedding season, weather",
    "→  Regional intelligence — South ≠ North ≠ Metro",
    "→  Oud demand spikes 4-6 wks before Eid;\n     system flags at 8 weeks",
]
for j, d in enumerate(demand_items):
    add_textbox(slide, 1.0, 2.7 + j * 0.5, 5, 0.5, d, font_size=13, color=TEXT_SECONDARY)

add_rounded_rect(slide, 1.5, 4.6, 2.5, 0.6, BG_CARD)
add_textbox(slide, 1.5, 4.6, 2.5, 0.3, '40%', font_size=24, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.9, 2.5, 0.3, 'Stockout reduction', font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Inventory Intelligence
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.2, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'INVENTORY INTELLIGENCE', font_size=11, color=BRAND_GOLD, bold=True)
inv_items = [
    "→  Computer vision: phone camera scans shelf,\n     counts units, flags gaps",
    "→  Auto-replenishment triggers to warehouse",
    "→  Dead stock alerts: 'This SKU hasn't moved\n     in 60 days at 14 stores'",
    "→  Franchise stock visibility — every shelf, every store",
]
for j, item in enumerate(inv_items):
    add_textbox(slide, 7.0, 2.7 + j * 0.5, 5.3, 0.5, item, font_size=13, color=TEXT_SECONDARY)

add_rounded_rect(slide, 8.0, 4.6, 2.5, 0.6, BG_CARD)
add_textbox(slide, 8.0, 4.6, 2.5, 0.3, '25%', font_size=24, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 8.0, 4.9, 2.5, 0.3, 'Overstock reduction', font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'Fragrance has brutal seasonality — each region peaks differently. The AI learns YOUR patterns.' Computer vision: 'Walk the store with your phone. 15 minutes — complete auditable inventory. No barcode scanning, no manual entry.' For franchise stores, inventory is usually a black box — this opens it up.")


# ============================================================
# SLIDE 9: STAFF PERFORMANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Your People Are Your Stores. Now You Can See.', font_size=36, color=WHITE, bold=True)

# Performance
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.5, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'PERFORMANCE INTELLIGENCE', font_size=11, color=BRAND_SECONDARY, bold=True)
perf_items = [
    "→  Individual scorecards: sales/hour, conversion,\n     basket size, upsell rate",
    "→  Store manager scorecards: team performance,\n     target attainment",
    "→  AI coaching: 'Ravi at Store 47 — high conversion\n     but low basket. Coach on upselling.'",
    "→  Gamification: leaderboards, streaks, benchmarking",
    "→  Area-level: staffing issues vs location issues",
]
for j, item in enumerate(perf_items):
    add_textbox(slide, 1.0, 2.7 + j * 0.55, 5, 0.55, item, font_size=13, color=TEXT_SECONDARY)

# Training
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.5, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'AI-POWERED TRAINING', font_size=11, color=BRAND_GOLD, bold=True)
train_items = [
    "→  Replaces travel-heavy trainer model entirely",
    "→  Interactive modules: product knowledge,\n     POS operation, selling techniques",
    "→  AI tests if associate can recommend correctly\n     before they face a customer",
    "→  Fragrance family quizzes, note profiles,\n     customer scenarios",
]
for j, item in enumerate(train_items):
    add_textbox(slide, 7.0, 2.7 + j * 0.55, 5.3, 0.55, item, font_size=13, color=TEXT_SECONDARY)

add_rounded_rect(slide, 8.0, 4.8, 3.5, 0.7, BG_CARD)
add_textbox(slide, 8.0, 4.8, 3.5, 0.35, '3 weeks → 5 days', font_size=22, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 8.0, 5.15, 3.5, 0.3, 'New hire ramp-up time', font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 6.3, 11, 0.4,
    'Dashboards are not surveillance — they make your best performers visible and your struggling performers get help earlier.',
    font_size=12, color=TEXT_MUTED)

add_notes(slide, "'100+ stores and growing = constant hiring. Every new hire is a revenue risk until trained. AI training compresses that window.' Performance dashboards feed the Command Center — area managers see staffing vs. location issues instantly. AI training eliminates the need for trainers to travel to every new store.")


# ============================================================
# SLIDE 10: L&D ACADEMY — GAMIFIED STAFF TRAINING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Bettroi Academy — Where Your Team Levels Up', font_size=34, color=WHITE, bold=True)

# Learning Platform
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.5, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'LEARNING PLATFORM', font_size=11, color=BRAND_SECONDARY, bold=True)
learn_items = [
    "→  Micro-learning: bite-sized modules on product\n     knowledge, selling, and brand story",
    "→  Role-based paths: new hire, store manager,\n     franchise partner, area head",
    "→  AI-personalized: curriculum adapts to each\n     associate's performance gaps",
    "→  Fragrance masterclasses: oud families, note\n     profiles, occasion pairing",
    "→  Multi-language: Hindi, English, regional —\n     learn in the language you think in",
    "→  Mobile-first: learn on the shop floor\n     between customers",
]
for j, item in enumerate(learn_items):
    add_textbox(slide, 1.0, 2.6 + j * 0.48, 5, 0.48, item, font_size=12, color=TEXT_SECONDARY)

# Gamification Engine
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.5, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'GAMIFICATION ENGINE', font_size=11, color=BRAND_GOLD, bold=True)
game_items = [
    "→  XP points for modules, quizzes, daily\n     challenges, and real sales milestones",
    "→  Leaderboards: store-level, region-level,\n     and national rankings",
    "→  Badges & certifications: 'Oud Expert',\n     'Top Seller', 'Customer Champion'",
    "→  Streak rewards: 7-day learning streaks\n     unlock perks and recognition",
    "→  Team battles: store vs. store competitions\n     tied to real KPIs",
    "→  Manager view: track team progress, identify\n     coaching opportunities",
]
for j, item in enumerate(game_items):
    add_textbox(slide, 7.0, 2.6 + j * 0.48, 5.3, 0.48, item, font_size=12, color=TEXT_SECONDARY)

# Stats
add_rounded_rect(slide, 0.8, 5.9, 3.7, 0.7, BG_CARD)
add_textbox(slide, 0.8, 5.9, 3.7, 0.35, '3x', font_size=24, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 6.25, 3.7, 0.3, 'Course completion rate vs\ntraditional training', font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 4.8, 5.9, 3.7, 0.7, BG_CARD)
add_textbox(slide, 4.8, 5.9, 3.7, 0.35, '60%', font_size=24, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 4.8, 6.25, 3.7, 0.3, 'Knowledge retention\nimprovement', font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 8.8, 5.9, 3.7, 0.7, BG_CARD)
add_textbox(slide, 8.8, 5.9, 3.7, 0.35, '85%', font_size=24, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 8.8, 6.25, 3.7, 0.3, 'Daily active engagement\nwith gamification', font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'Traditional retail training is broken — send a trainer to 100 stores, or fly staff to HQ. Neither scales. Bettroi Academy makes every phone a training center.' Gamification is not gimmicks — it drives daily engagement. Staff who complete learning paths sell 22% more. The leaderboards create healthy competition — top performers get visibility, struggling staff get support. Certifications like 'Oud Expert' give associates pride and customers confidence.")


# ============================================================
# SLIDE 11: CUSTOMER & WHATSAPP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Know Your Customer. Reach Them. Keep Them.', font_size=36, color=WHITE, bold=True)

# Customer Analytics
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.2, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'CUSTOMER ANALYTICS', font_size=11, color=BRAND_SECONDARY, bold=True)
cust_items = [
    "→  Unified profiles: retail, Shoppers Stop, Nykaa,\n     Amazon, airports",
    "→  Segments: first-time, repeat, lapsed (90+ days),\n     high-value",
    "→  'Buys new EDP every 3 months — trigger\n     reminder at 2.5 months'",
    "→  LTV prediction: 'These 2,000 customers =\n     40% of next year's revenue'",
    "→  Churn alerts: '142 high-value customers\n     silent for 60+ days'",
]
for j, item in enumerate(cust_items):
    add_textbox(slide, 1.0, 2.6 + j * 0.52, 5, 0.52, item, font_size=12, color=TEXT_SECONDARY)

# WhatsApp CRM
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.2, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'WHATSAPP CRM & CLIENTELING', font_size=11, color=BRAND_GOLD, bold=True)
wa_items = [
    "→  Auto messages: order confirm, restock\n     reminders, birthday offers",
    "→  Campaigns: 'Send Eid collection preview to\n     all oud buyers from last Eid'",
    "→  Clienteling: top 20 customers with prefs\n     & next-buy suggestions",
    "→  Two-way chat: questions, repeat orders\n     via WhatsApp",
]
for j, item in enumerate(wa_items):
    add_textbox(slide, 7.0, 2.6 + j * 0.52, 5.3, 0.52, item, font_size=12, color=TEXT_SECONDARY)

# Example
add_rounded_rect(slide, 6.8, 4.8, 5.7, 1.1, RGBColor(0x1A, 0x18, 0x10), BRAND_GOLD)
add_textbox(slide, 7.0, 4.85, 5.3, 0.3, 'Customer buys 50ml EDP on Jan 15. On Mar 25:', font_size=11, color=TEXT_MUTED)
add_textbox(slide, 7.0, 5.2, 5.3, 0.5,
    '"Hi Priya, running low on Wisal? Reorder now + explore our spring collection. Reply YES to order."',
    font_size=12, color=TEXT_PRIMARY, bold=True)

add_notes(slide, "'In fragrance, replenishment is predictable. 50ml EDP lasts 2-3 months.' Mention all channels naturally — retail, Shoppers Stop, Nykaa, Amazon, airports. '97% of customers are on WhatsApp. Manual WhatsApp CRM is chaos. This automates the right message, right customer, right time.'")


# ============================================================
# SLIDE 11: DYNAMIC PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'AI MODULE', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'The Right Price. The Right Store. The Right Time.', font_size=36, color=WHITE, bold=True)

# Intelligent Pricing
add_rounded_rect(slide, 0.8, 2.2, 5.5, 2.8, BG_CARD)
add_textbox(slide, 1.0, 2.3, 4, 0.3, 'INTELLIGENT PRICING', font_size=11, color=BRAND_SECONDARY, bold=True)
price_items = [
    "→  AI-recommended adjustments based on demand,\n     stock levels, competition",
    "→  Markdown: '15% off moves volume; 30% destroys\n     margin — AI finds the sweet spot'",
    "→  Festive: 'Eid gift sets — last year 20% off moved\n     3x volume, test 15% this year'",
]
for j, item in enumerate(price_items):
    add_textbox(slide, 1.0, 2.8 + j * 0.65, 5, 0.65, item, font_size=13, color=TEXT_SECONDARY)

# Channel Governance
add_rounded_rect(slide, 6.8, 2.2, 5.7, 2.8, BG_CARD)
add_textbox(slide, 7.0, 2.3, 4, 0.3, 'CHANNEL GOVERNANCE', font_size=11, color=BRAND_GOLD, bold=True)
gov_items = [
    "→  Channel harmonization: retail = Nykaa = Amazon",
    "→  Franchise pricing governance: floors & ceilings\n     with controlled flexibility",
    "→  Alert on marketplace price drift",
]
for j, item in enumerate(gov_items):
    add_textbox(slide, 7.0, 2.8 + j * 0.65, 5.3, 0.65, item, font_size=13, color=TEXT_SECONDARY)

# Stats
add_rounded_rect(slide, 0.8, 5.3, 5.5, 1.0, BG_CARD)
add_textbox(slide, 0.8, 5.35, 5.5, 0.5, '3-8%', font_size=36, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 5.85, 5.5, 0.3, 'Margin improvement', font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 6.8, 5.3, 5.7, 1.0, BG_CARD)
add_textbox(slide, 6.8, 5.35, 5.7, 0.5, 'Zero', font_size=36, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 6.8, 5.85, 5.7, 0.3, 'Channel price conflicts', font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 6.6, 11, 0.4,
    'Most retailers over-discount because they lack price elasticity data. The AI finds the discount that moves inventory without destroying margin.',
    font_size=12, color=TEXT_MUTED)

add_notes(slide, "'If Nykaa price < store price, you train customers to buy online. If franchise ≠ company-owned, brand confusion.' This gets the CFO's attention. Reinforces the need for clean POS data. Markdown optimization alone typically pays for the platform.")


# ============================================================
# SLIDE 12: AI ADVANTAGE SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'TECHNOLOGY', font_size=11, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, '9 AI Engines. One Platform. Compounding Intelligence.', font_size=34, color=WHITE, bold=True)

ai_engines = [
    ("🧴", "FragranceAI", "Personalized scent recommendations from purchase + preference data", BRAND_PRIMARY),
    ("📈", "DemandAI", "Predictive demand forecasting with cultural event awareness", BRAND_SECONDARY),
    ("👁", "VisionAI", "Computer vision inventory counting & shelf compliance", BRAND_GOLD),
    ("💰", "PricingAI", "Dynamic pricing optimization across channels & store types", BRAND_ACCENT),
    ("👥", "PeopleAI", "Staff performance scoring, coaching triggers & scheduling", BRAND_PRIMARY),
    ("🎓", "AcademyAI", "Gamified L&D platform with adaptive learning paths & certifications", BRAND_SECONDARY),
    ("🎯", "CustomerAI", "Segmentation, lifetime value prediction & churn prevention", BRAND_GOLD),
    ("💬", "ConversationAI", "WhatsApp automation, chatbot & clienteling prompts", BRAND_ACCENT),
    ("🔍", "InsightAI", "'What were my top 5 stores this Eid?' — instant answers", BRAND_PRIMARY),
]

for i, (icon, name, desc, color) in enumerate(ai_engines):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.1 + row * 1.1
    add_rounded_rect(slide, x, y, 3.8, 0.9, BG_CARD)
    add_textbox(slide, x + 0.15, y + 0.05, 0.5, 0.5, icon, font_size=18)
    add_textbox(slide, x + 0.7, y + 0.08, 2.8, 0.3, name, font_size=13, color=color, bold=True)
    add_textbox(slide, x + 0.7, y + 0.4, 2.9, 0.4, desc, font_size=10, color=TEXT_SECONDARY)

add_textbox(slide, 0.8, 6.5, 11, 0.5,
    'Each AI engine gets smarter over time. And they feed each other — DemandAI uses PricingAI data, FragranceAI uses CustomerAI profiles, InsightAI sits on top of everything.',
    font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'Any vendor can say AI-powered. We show exactly which models power which capabilities.' InsightAI example: 'Imagine your CEO in Dubai typing: What were my top 10 stores this Eid vs last Eid? — instant answer with chart. No analyst needed.' This is the second 'Dubai' reference — keep it casual.")


# ============================================================
# SLIDE 13: ROI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'BUSINESS CASE', font_size=11, color=BRAND_SECONDARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'What the Numbers Look Like', font_size=36, color=WHITE, bold=True)

# Operational Savings
add_rounded_rect(slide, 0.8, 2.1, 5.5, 3.2, BG_CARD)
add_textbox(slide, 1.0, 2.2, 4, 0.3, 'OPERATIONAL SAVINGS', font_size=11, color=BRAND_SECONDARY, bold=True)
savings = [
    ("15-20 hrs/store", "Monthly billing time saved"),
    ("30+ hrs/month", "HQ manual dashboard work eliminated"),
    ("70% reduction", "Inventory counting time"),
    ("3 wks → 5 days", "New hire training time"),
    ("Near-zero", "Promo / pricing errors"),
]
for j, (val, label) in enumerate(savings):
    add_textbox(slide, 1.0, 2.7 + j * 0.48, 2.2, 0.3, val, font_size=13, color=BRAND_SECONDARY, bold=True)
    add_textbox(slide, 3.3, 2.7 + j * 0.48, 2.8, 0.3, label, font_size=12, color=TEXT_SECONDARY)

# Revenue Uplift
add_rounded_rect(slide, 6.8, 2.1, 5.7, 3.2, BG_CARD)
add_textbox(slide, 7.0, 2.2, 4, 0.3, 'REVENUE UPLIFT', font_size=11, color=BRAND_GOLD, bold=True)
uplift = [
    ("22-35%", "Basket size increase (AI recs)"),
    ("40%", "Stockout reduction = captured sales"),
    ("15-20%", "Repeat purchase improvement (WhatsApp)"),
    ("3-8%", "Margin improvement (dynamic pricing)"),
]
for j, (val, label) in enumerate(uplift):
    add_textbox(slide, 7.0, 2.7 + j * 0.52, 2.0, 0.3, val, font_size=14, color=BRAND_GOLD, bold=True)
    add_textbox(slide, 9.2, 2.7 + j * 0.52, 3.0, 0.3, label, font_size=12, color=TEXT_SECONDARY)

# Bottom CTA
add_rounded_rect(slide, 0.8, 5.6, 11.7, 1.3, RGBColor(0x14, 0x12, 0x28), BRAND_PRIMARY)
add_textbox(slide, 0.8, 5.65, 11.7, 0.3, 'For a 100-store retail chain', font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 5.95, 11.7, 0.5, '₹8-12 Cr Annual Impact', font_size=36, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 6.45, 11.7, 0.3, 'Within 18 months of full deployment  |  Platform payback: 4-6 months',
    font_size=13, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'The operational savings alone — POS and dashboards — typically pay for the platform within 4-6 months. Everything else is upside.' The '100-store chain' language is generic — let them do the math. Estimate is conservative — doesn't account for peak season stockout recovery.")


# ============================================================
# SLIDE 14: IMPLEMENTATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'ROLLOUT', font_size=11, color=BRAND_SECONDARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Value in 8 Weeks. Full Platform in 6 Months.', font_size=36, color=WHITE, bold=True)

# Phase 1
add_rounded_rect(slide, 0.8, 2.2, 11.7, 1.2, RGBColor(0x1E, 0x16, 0x0A), BRAND_GOLD)
add_textbox(slide, 1.0, 2.25, 0.6, 0.5, 'P1', font_size=22, color=BRAND_GOLD, bold=True)
add_textbox(slide, 1.8, 2.25, 6, 0.35, 'Foundation — Smart POS + Command Center', font_size=16, color=WHITE, bold=True)
add_textbox(slide, 9.5, 2.25, 2.5, 0.35, 'Weeks 1-8', font_size=13, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.RIGHT)
add_textbox(slide, 1.8, 2.7, 10, 0.5, 'Deploy across 10-15 pilot stores. HQ dashboards live. Staff trained. Parallel running with existing systems — zero disruption.',
    font_size=12, color=TEXT_SECONDARY)

# Phase 2
add_rounded_rect(slide, 0.8, 3.6, 11.7, 1.2, RGBColor(0x0E, 0x1A, 0x18), BRAND_SECONDARY)
add_textbox(slide, 1.0, 3.65, 0.6, 0.5, 'P2', font_size=22, color=BRAND_SECONDARY, bold=True)
add_textbox(slide, 1.8, 3.65, 6, 0.35, 'Intelligence — Forecasting, Inventory AI, Staff Performance', font_size=16, color=WHITE, bold=True)
add_textbox(slide, 9.5, 3.65, 2.5, 0.35, 'Weeks 9-16', font_size=13, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.RIGHT)
add_textbox(slide, 1.8, 4.1, 10, 0.5, 'Activate demand forecasting and inventory intelligence. Staff performance AI and training modules go live. Customer data unification begins.',
    font_size=12, color=TEXT_SECONDARY)

# Phase 3
add_rounded_rect(slide, 0.8, 5.0, 11.7, 1.2, RGBColor(0x14, 0x12, 0x28), BRAND_PRIMARY)
add_textbox(slide, 1.0, 5.05, 0.6, 0.5, 'P3', font_size=22, color=BRAND_PRIMARY, bold=True)
add_textbox(slide, 1.8, 5.05, 6, 0.35, 'Experience — FragranceAI, WhatsApp CRM, Dynamic Pricing', font_size=16, color=WHITE, bold=True)
add_textbox(slide, 9.5, 5.05, 2.5, 0.35, 'Weeks 17-24', font_size=13, color=BRAND_SECONDARY, bold=True, alignment=PP_ALIGN.RIGHT)
add_textbox(slide, 1.8, 5.5, 10, 0.5, 'Full AI suite deployed. WhatsApp CRM live. Dynamic pricing activated. Rollout across all stores.',
    font_size=12, color=TEXT_SECONDARY)

# Bottom cards
cards_impl = [
    ("Dedicated Team", "Not a self-serve setup"),
    ("Parallel Running", "Zero disruption to operations"),
    ("Franchise Playbook", "Onboarding guide included"),
]
for i, (title, sub) in enumerate(cards_impl):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 6.4, 3.7, 0.8, BG_CARD)
    add_textbox(slide, x, 6.45, 3.7, 0.35, title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 6.75, 3.7, 0.3, sub, font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'We always start with the foundation — POS and dashboards — because they generate the data for everything else.' Pilot: 'See results at 10-15 stores before committing. We earn the expansion.' Phase 1 focus is architecturally justified: 'You cannot run AI on bad data.'")


# ============================================================
# SLIDE 15: WHY BETTROI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.4, 'WHY US, WHY NOW', font_size=11, color=BRAND_SECONDARY, bold=True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, 'Built for This Moment', font_size=40, color=WHITE, bold=True)

pillars = [
    ("Retail-Native", "Built exclusively for retail. Not horizontal SaaS adapted — retail-native, retail-only. We understand the floor.", BRAND_PRIMARY),
    ("India-First", "GST complexity, UPI payments, Hindi & regional language support, offline-first architecture, WhatsApp integration.", BRAND_SECONDARY),
    ("AI-Native", "AI is not a feature we bolted on. It is the architecture. Every module generates data, every data point trains the AI.", BRAND_GOLD),
    ("Scale-Ready", "Designed for 100+ store chains. Multi-format, multi-channel, franchise-ready. Scales from pilot to national.", BRAND_ACCENT),
]
for i, (title, desc, color) in enumerate(pillars):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 2.1 + row * 1.6
    add_rounded_rect(slide, x, y, 5.8, 1.4, BG_CARD, color)
    add_textbox(slide, x + 0.3, y + 0.15, 5, 0.35, title, font_size=18, color=color, bold=True)
    add_textbox(slide, x + 0.3, y + 0.55, 5.2, 0.7, desc, font_size=12, color=TEXT_SECONDARY)

# CTA
add_rounded_rect(slide, 2.5, 5.5, 8.3, 1.6, RGBColor(0x14, 0x12, 0x28), BRAND_PRIMARY)
add_textbox(slide, 2.5, 5.6, 8.3, 0.5, "Let's Start With a Pilot", font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.5, 6.1, 8.3, 0.35, '10-15 stores  ·  8 weeks  ·  Let the results speak', font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.5, 6.6, 8.3, 0.35, 'BETTROI  |  hello@bettroi.com', font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_notes(slide, "'The retail chains that modernize this year will define the next decade.' 'We're not asking you to be a guinea pig. We're asking you to be a leader.' End with: '10-15 stores, 8 weeks, let the results speak.' Leave contact info on screen during Q&A.")


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), 'Bettroi_Retail_Intelligence_Platform.pptx')
prs.save(output_path)
print(f"[OK] PowerPoint saved to: {output_path}")
print(f"  Slides: {len(prs.slides)}")
